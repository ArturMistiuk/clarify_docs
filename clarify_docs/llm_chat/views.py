import os
import tempfile
import docx2txt
from PyPDF2 import PdfReader
from django.contrib.auth.decorators import login_required
from langchain.callbacks import get_openai_callback

from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import Chroma
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
# Create your views here.
from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse
import openai
from dotenv import load_dotenv
from langchain.vectorstores.faiss import FAISS
from pptx import Presentation

from .forms import PDFUploadForm, PDFUpdateForm, PDFDocumentForm2
from .models import ChatMessage, PDFDocument, UserData

openai.api_key = "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"
CHUNK_SIZE = 1024
CHUNK_OVERLAP = 200
MAX_LEN = 512
TEMPERATURE = 0.5
load_dotenv()


def main(request):
    # avatar = Avatar.objects.filter(user_id=request.user.id).first()
    return render(request, 'llm_chat/index.html', context={}) #home.html


def get_pdf_text(file):
    text = None
    if file:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(file.read())
            temp_file.flush()
            if file.name.endswith('.pdf'):
                pdf_reader = PdfReader(temp_file.name)
                text = ''.join(page.extract_text() for page in pdf_reader.pages)
            elif file.name.endswith('.txt'):
                with open(temp_file.name, 'r') as f:
                    text = f.read()
            elif file.name.endswith('.docx'):
                text = docx2txt.process(temp_file.name)
            elif file.name.endswith('.pptx'):
                prs = Presentation(temp_file.name)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = '\n'.join(text_runs)
    return text


def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    api_key = os.getenv("OPENAI_API_KEY")
    embeddings = OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=api_key)

    # embeddings = OllamaEmbeddings(model_name=EMBEDDING_MODEL)
    # embeddings = HuggingFaceInstructEmbeddings(
    #     model_name="hkunlp/instructor-large", model_kwargs={"device": 'cpu'}
    # )
    # embeddings = HuggingFaceInstructEmbeddings(model_name=EMBEDDING_MODEL)
    vectorstore = FAISS.from_texts(text_chunks, embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(model="gpt-3.5-turbo")
    # llm = AutoModelForCausalLM.from_pretrained("TheBloke/Llama-2-7b-Chat-GGUF",
    #                                            model_file="llama-2-7b-chat.q4_K_M.gguf", model_type="llama",
    #                                            gpu_layers=0, from_tf=True)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


# def home(request):
#     try:
#         # if the session does not have a messages key, create one
#         if 'messages' not in request.session:
#             request.session['messages'] = [
#                 {"role": "system",
#                  "content": "You are now chatting with a user, provide them with comprehensive, short and concise answers."},
#             ]
#
#         if request.method == 'POST':
#             # get the prompt from the form
#             prompt = request.POST.get('prompt')
#             # get the temperature from the form
#             temperature = float(request.POST.get('temperature', 0.1))
#             # append the prompt to the messages list
#             request.session['messages'].append({"role": "user", "content": prompt})
#             # set the session as modified
#             request.session.modified = True
#             # call the openai API
#             response = openai.ChatCompletion.create(
#                 model="gpt-3.5-turbo",
#                 messages=request.session['messages'],
#                 temperature=temperature,
#                 max_tokens=1000,
#             )
#             # format the response
#             formatted_response = response['choices'][0]['message']['content']
#             # append the response to the messages list
#             request.session['messages'].append({"role": "llm_chat", "content": formatted_response})
#             request.session.modified = True
#             # redirect to the home page
#             context = {
#                 'messages': request.session['messages'],
#                 'prompt': '',
#                 'temperature': temperature,
#             }
#             return render(request, 'llm_chat/home.html', context)
#         else:
#             # if the request is not a POST request, render the home page
#             context = {
#                 'messages': request.session['messages'],
#                 'prompt': '',
#                 'temperature': 0.1,
#             }
#             return render(request, 'llm_chat/home.html', context)
#     except Exception as e:
#         print(e)
#         # if there is an error, redirect to the error handler
#         return redirect('error_handler')
@login_required(login_url="/login/")
def upload_pdf(request):
    user = request.user

    try:
        user_data = UserData.objects.get(user=user)
    except UserData.DoesNotExist:
        user_data = UserData.objects.create(user=user, total_files_uploaded=0, total_questions_asked=0)

    if request.method == 'POST': #and request.FILES.get('pdf_document'):
        form = PDFUploadForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_document = request.FILES['pdf_document']

            # Перевірка розширення файлу
            _, file_extension = os.path.splitext(pdf_document.name)
            if file_extension.lower() not in [".pdf", ".txt", ".docx", ".pptx"]:
                return JsonResponse({'error': 'Only PDF, TXT, DOCX, PPTX files'}, status=400)

            if pdf_document.size > 50 * 1024 * 1024:  # Розмір файлу понад 50 МБ
                return JsonResponse({'error': "File size exceeds 50 MB."}, status=400)

            # Check if a file with the same name already exists for this user
            if PDFDocument.objects.filter(user=user, title=pdf_document.name).exists():
                return JsonResponse({'error': 'A file with the same name already exists.'}, status=400)

            # Збільшити кількість завантажених файлів користувача
            user_data.total_files_uploaded += 1
            user_data.save()

            # Зберегти PDF-документ у базі даних
            pdf = PDFDocument(user=user, title=pdf_document.name)
            pdf.documentContent = get_pdf_text(pdf_document)
            pdf.save()
            return JsonResponse({'success': 'File uploaded successfully.'})

    else:
        form = PDFUploadForm()
    user_pdfs = PDFDocument.objects.filter(user=request.user)
    # not needed for now
    chat_message = ChatMessage.objects.all()
    return render(request, 'llm_chat/home.html', {'form': form, 'user_pdfs': user_pdfs})


@login_required(login_url="/login/")
def ask_question(request):
    user = request.user

    try:
        user_data = UserData.objects.get(user=user)
    except UserData.DoesNotExist:
        # Обробка ситуації, коли відсутній запис UserData для користувача.
        # Можна створити запис за замовчуванням тут.
        user_data = UserData.objects.create(user=user, total_files_uploaded=0, total_questions_asked=0)

    chat_history = ChatMessage.objects.filter(user=request.user).order_by('timestamp')[:10]
    chat_response = ''
    user_pdfs = PDFDocument.objects.filter(user=request.user)
    user_question = ""
    selected_pdf = None  # Змінено з selected_pdf_id на об'єкт PDFDocument

    if request.method == 'POST':
        # Check if the user has exceeded the limit for questions per file
        user_question = request.POST.get('user_question')
        selected_pdf_id = request.POST.get('selected_pdf')
        selected_pdf = get_object_or_404(PDFDocument, id=selected_pdf_id)
        text_chunks = get_text_chunks(selected_pdf.documentContent)

        knowledge_base = get_vectorstore(text_chunks)
        conversation_chain = get_conversation_chain(knowledge_base)

        with get_openai_callback() as cb:
            response = conversation_chain({'question': user_question})

        chat_response = response["answer"]
        chat_message = ChatMessage(user=request.user, message=user_question, answer=chat_response,
                                   pdf_document=selected_pdf)  # Передаємо об'єкт PDFDocument
        user_data.total_questions_asked += 1
        user_data.save()
        chat_message.save()

    # Отримуємо повідомлення, які відносяться до обраного PDFDocument
    chat_message = ChatMessage.objects.filter(user=request.user, pdf_document=selected_pdf).order_by('timestamp')

    context = {'chat_response': chat_response, 'chat_history': chat_history, 'user_question': user_question,
               'user_pdfs': user_pdfs, 'chat_message': chat_message}

    return render(request, 'llm_chat/home.html', context)

@login_required(login_url="/login/")
def get_chat_history(request):
    pdf_id = request.GET.get('pdf_id')
    if not pdf_id:
        return JsonResponse({"error": "PDF ID not provided"}, status=400)

    # Тут ви отримуєте історію чату для даного pdf_id
    # Це припущення, що у вас є модель ChatMessage або щось подібне
    messages = ChatMessage.objects.filter(pdf_document_id=pdf_id).values('message', 'answer', 'timestamp')

    return JsonResponse(list(messages), safe=False)


def new_chat(request):
    # clear the messages list
    request.session.pop('messages', None)
    return redirect('home')


# this is the view for handling errors
def error_handler(request):
    return render(request, 'llm_chat/404.html')

# def get_completion(prompt):
#     print(prompt)
#     query = openai.Completion.create(
#         engine="text-davinci-003",
#         prompt=prompt,
#         max_tokens=512,
#         n=1,
#         stop=None,
#         temperature=0.5,
#     )
#
#     response = query.choices[0].text
#     print(response)
#     return response
# llm = ChatOpenAI()
# embeddings = OpenAIEmbeddings()
# def query_view(request):
#
#     if request.method == 'POST':
#         prompt = request.POST.get('prompt')
#
#         # user_question = request.POST.get('user_question')
#         selected_pdf = request.POST.get('pdf_content')
#         # selected_pdf = get_object_or_404(PDFDocument, id=selected_pdf_id)
#         text_chunks = get_text_chunks(selected_pdf)
#
#         knowledge_base = get_vectorstore(text_chunks)
#         conversation_chain = get_conversation_chain(knowledge_base)
#
#         chat_response = conversation_chain({'question': prompt})
#
#         response = chat_response["answer"]
#
#
#
#
#         # response = get_completion(prompt)
#         return JsonResponse({'response': response})
#     return render(request, 'llm_chat/index.html')
#
#
# def index(request):
#     pdf_content = None
#
#     if request.method == 'POST' and request.FILES.get('pdf_file'):
#         pdf_file = request.FILES['pdf_file']
#
#         # Process the PDF content
#         pdf_content = process_pdf(pdf_file)
#
#     return render(request, 'index.html', {'pdf_content': pdf_content})
#
#
# def process_pdf(pdf_file):
#     pdf_text = ''
#     try:
#         # Read the uploaded PDF file
#         pdf = PdfReader(pdf_file)
#
#         # Extract text from each page
#         for page in pdf.pages:
#             pdf_text += page.extract_text()
#     except Exception as e:
#         pdf_text = f"Error processing PDF: {str(e)}"
#
#     return pdf_text
#
# def get_text_chunks(text):
#     text_splitter = RecursiveCharacterTextSplitter(
#         # separator="\n",
#         chunk_size=CHUNK_SIZE,
#         chunk_overlap=CHUNK_OVERLAP,
#         length_function=len
#     )
#     chunks = text_splitter.split_text(text)
#     return chunks
#
#
# def get_vectorstore(text_chunks):
#     # embeddings = OpenAIEmbeddings()
#     # embeddings = HuggingFaceInstructEmbeddings(model_name=EMBEDDING_MODEL)
#     vectorstore = Chroma.from_texts(texts=text_chunks, embedding=embeddings, persist_directory="vectorstore")
#     return vectorstore
#
#
# def get_conversation_chain(vectorstore):
#     # llm = ChatOpenAI()
#     # llm = HuggingFaceHub(repo_id=LLM_MODEL, model_kwargs={"temperature": TEMPERATURE, "max_length": MAX_LEN})
#     # llm = CTransformers(model='llama-2-7b-chat.Q4_K_M.gguf',
#     #                     model_type='llama',
#     #                     temperature=0.75,
#     #                     max_tokens=2000,
#     #                     top_p=1,
#     #                     verbose=True,  # Verbose is required to pass to the callback manager
#     #                     )
#     memory = ConversationBufferMemory(
#         memory_key='chat_history', return_messages=True)
#     conversation_chain = ConversationalRetrievalChain.from_llm(
#         llm=llm,
#         retriever=vectorstore.as_retriever(),
#         memory=memory
#     )
#     return conversation_chain
